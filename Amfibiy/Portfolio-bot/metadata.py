"""Извлечение и нормализация метаданных работ: группа/автор/предмет/тип работы и т.п."""
import os
import re
import asyncio
import cv2
import numpy as np
import io
import json
import uuid
import shutil
import logging
import subprocess
import tempfile
import docx
from pypdf import PdfReader
import pikepdf
import img2pdf
from pdf2image import convert_from_path
from google import genai
from google.genai import types
from rapidfuzz import fuzz
from pptx import Presentation
import fitz
from PIL import Image

ALLOWED_WORK_TYPES = ["ЛР", "ПЗ", "КР", "КП", "ИЗ", "СРС", "ДКР", "СР", "ПР", "УП", "ПП"]
from text_helpers import clean_text_for_parsing, normalize_text
from ocr import extract_raw_text_gemini

def get_subject_abbreviation(text: str) -> str:
    text = normalize_text(text)
    if not text:
        return "Предмет"
    if " " not in text:
        return text
    if text.isupper():
        return text
    words = re.findall(r'[а-яА-ЯёЁa-zA-Z0-9]+', text)
    if not words:
        return "Предмет"
    prepositions = {"и", "в", "на", "с", "по", "для", "о", "об", "к", "из"}
    ignore = {"тип", "работы", "работа", "лабораторная", "практическая",
              "отчет", "отчёт", "дисциплина", "дисциплине"}
    abbr = ""
    for w in words:
        if w.lower() in ignore:
            continue
        if w.isupper() and len(w) > 1:
            abbr += w
        elif w.lower() in prepositions:
            abbr += w[0].lower()
        else:
            abbr += w[0].upper()
    return abbr if abbr else "Предмет"

def extract_authors(text: str) -> list:
    authors = []
    start_pattern = r'(?:Выполнили\s+студенты|Выполнил\s+студент|Выполнили|Выполнил|Авторы|Автор|Студенты|Студент)[\s:]+'
    match = re.search(start_pattern, text, re.IGNORECASE)
    if not match:
        return authors

    block = text[match.end():]
    end_match = re.search(
        r'(Руководитель|Проверил|Преподаватель|Принял|Оценил|Екатеринбург|г \. |202\d|Отзыв|Содержание|Введение)',
        block, re.IGNORECASE
    )
    if end_match:
        block = block[:end_match.start()]

    stop_words = {'Группа', 'Студент', 'Студенты', 'Курс', 'Семестр', 'Работу', 'Проект', 'Гр'}
    pattern = r'([А-ЯЁ][а-яё]{2,})[\s,]+([А-ЯЁ])[\.\s,]*([А-ЯЁ])?[\.\s,]?'
    for m in re.finditer(pattern, block):
        if m.group(1).title() in stop_words:
            continue
        author = f"{m.group(1)} {m.group(2)}{m.group(3) or ''}".replace('.', '').strip()
        if author not in authors:
            authors.append(author)

    if not authors:
        for m in re.finditer(r'([А-ЯЁ][а-яё]{2,})\s+([А-ЯЁ][а-яё]{2,})', block):
            if m.group(1).title() not in stop_words and m.group(2).title() not in stop_words:
                author = f"{m.group(1)} {m.group(2)}"
                if author not in authors:
                    authors.append(author)
    return authors

def extract_department(text: str) -> str:
    """Извлекает код кафедры из текста, приводя к одному из трёх вариантов."""
    if not text:
        return ""
    text_up = text.upper()
    if "ИНФОКОММУНИКАЦИОННЫХ ТЕХНОЛОГИЙ" in text_up or "ИТИМС" in text_up:
        return "ИТиМС"
    if "МНОГОКАНАЛЬНОЙ ЭЛЕКТРИЧЕСКОЙ СВЯЗИ" in text_up or "МЭС" in text_up:
        return "МЭС"
    if "ИНФОРМАЦИОННЫХ СИСТЕМ" in text_up or "ИСТ" in text_up:
        return "ИСТ"
    # Альтернативные написания
    if "ИТИМС" in text_up:
        return "ИТиМС"
    if "МЭС" in text_up:
        return "МЭС"
    if "ИСТ" in text_up:
        return "ИСТ"
    return ""

def clean_group_name(group_str: str) -> str:
    if not group_str:
        return "БЕЗ_ГРУППЫ"
    
    clean = re.sub(r'(гр\.?|группа|group)\s*', '', group_str, flags=re.IGNORECASE)
    clean = clean.upper().strip().replace(" ", "")
    
    clean = re.sub(r'[^А-ЯЁA-Z0-9\-]', '', clean)
    
    if len(clean) >= 5 and clean[-1] == '6' and clean[-2].isdigit():
        clean = clean[:-1] + 'Б'
    
    match = re.match(r'^([А-ЯЁA-Z]{2,4})(\d{2,4})([А-ЯЁA-Z]?)$', clean)
    if match:
        letters, digits, letter = match.groups()
        if '-' not in clean:
            clean = f"{letters}-{digits}{letter}"
    
    return clean or "БЕЗ_ГРУППЫ"

def clean_subject_name(subject_str: str) -> str:
    if not subject_str:
        return "Предмет"
    clean = re.sub(r'[\u2010\u2011\u2012\u2013\u2014\u2015\u2212]', '-', subject_str)
    clean = re.sub(r'\s*-\s*', '-', clean)
    clean = clean.replace(" ", "-").replace("_", "-")
    return re.sub(r'-+', '-', clean)

def clean_author_name(author_str) -> str:
    if isinstance(author_str, (list, tuple, set)):
        author_str = next((a for a in author_str if a), None)
    if not author_str or not isinstance(author_str, str):
        return "Автор"
    clean = author_str.replace(".", "")
    clean = re.sub(r'\s([А-ЯЁа-яё])\s([А-ЯЁа-яё])$', r' \1\2', clean)
    return " ".join(clean.split())

def normalize_author_name(name: str) -> str:
    return clean_author_name(name)

def generate_safe_filename(data: dict, specific_author: str = None) -> str:
    group = clean_group_name(data.get("group"))
    if specific_author:
        author_str = specific_author
    else:
        authors = data.get("authors")
        author_str = (authors[0] if authors and isinstance(authors, list) else None) \
                     or data.get("author") or "Автор"
    author_str = clean_author_name(author_str)
    subject = clean_subject_name(data.get("subject"))
    work_type = data.get("work_type") or ""
    number = data.get("work_number") or ""
    type_part = f"_{work_type}{number}" if work_type and work_type != "РАБ" else (f"_{number}" if number else "")
    filename = f"{group}_{author_str}_{subject}{type_part}.pdf"
    return re.sub(r'[\\/*?:"<>|]', '', filename)

def check_filename_format(filename: str) -> dict:
    result = {"valid": False, "has_pdf": False, "parsed": None}
    name = os.path.basename(filename)
    if not name.lower().endswith('.pdf'):
        return result
    result["has_pdf"] = True
    pattern = re.compile(
        r'^([А-ЯЁA-Z0-9][А-ЯЁA-Z0-9\-]+)_'
        r'([А-ЯЁ][а-яё]+ [А-ЯЁ]{1,2}|[А-ЯЁ][а-яё]+[А-ЯЁ]{2})_'
        r'([А-ЯЁA-Zа-яёa-z0-9\-]+)'
        r'(?:_(ЛР|ПЗ|КР|КП|ИЗ|СРС|ДКР|СР|ПР|УП|ПП)\d*)?'
        r'\.pdf$',
        re.IGNORECASE
    )
    m = pattern.match(name)
    if m:
        result["valid"] = True
        result["parsed"] = {
            "group": m.group(1),
            "author": m.group(2),
            "subject": m.group(3),
            "work_type": m.group(4) or "",
        }
    return result

def get_work_type_from_text(text: str) -> str:
    text_lower = text.lower()
    mapping = {
        "курсовой проект": "КП", "курсовая работа": "КР",
        "лабораторная работа": "ЛР", "лабораторная": "ЛР",
        "практическая работа": "ПЗ", "практическая": "ПЗ",
        "индивидуальное задание": "ИЗ", "индивидуальная работа": "ИЗ",
        "учебная практика": "УП", "производственная практика": "ПП",
    }
    for phrase, code in mapping.items():
        if phrase in text_lower:
            return code
    return None

def parse_gemini_ocr_text(text: str) -> dict:
    text = clean_text_for_parsing(text)
    group = author = subject = work_type = num = None
    subj_text = "Предмет"

    m = re.search(r'Группа:\s*([^\n]+)', text, re.IGNORECASE)
    if m:
        group = m.group(1).strip().upper().replace(" ", "-")
        if len(group) >= 5 and group[-1] == '6' and group[-2].isdigit():
            group = group[:-1] + 'Б'

    authors = extract_authors(text)
    if authors:
        author = authors[0]
    else:
        m2 = re.search(r'Выполнил студент:\s*([^\n]+)', text, re.IGNORECASE)
        if m2:
            av = m2.group(1).strip()
            am = re.search(r'([А-ЯЁ][а-яё]+)\s+([А-ЯЁ])[\.\s,]*([А-ЯЁ])?', av)
            if am:
                author = f"{am.group(1)} {am.group(2)}{am.group(3) or ''}".replace('.', '')
            else:
                author = av.replace('.', '')

    m3 = re.search(r'По дисциплине:\s*([^\n]+)', text, re.IGNORECASE)
    if m3:
        subj_text = re.sub(r'^[«"\'\(]+|[»"\'\)]+$', '', m3.group(1).strip())
        subject = get_subject_abbreviation(subj_text)

    type_map = {
        "лабораторн": "ЛР", "лаб. раб": "ЛР", 
        "практическ": "ПЗ", "пр. раб": "ПЗ", "практ": "ПЗ",
        "курсовой проект": "КП", "курсовая проект": "КП",
        "курсовая работа": "КР", "курсовой работа": "КР", "курсов": "КР", 
        "индивидуальн": "ИЗ",
        "учебная практика": "УП", "производственная практика": "ПП",
    }
    m4 = re.search(r'Тип работы:\s*([^\n]+)', text, re.IGNORECASE)
    if m4:
        tv = m4.group(1).lower()
        for key, val in type_map.items():
            if key in tv:
                work_type = val
                break
        if not work_type:
            alt = re.search(r'\b(ЛР|ПЗ|КР|КП|ИЗ|СРС|ДКР|СР|ПР|УП|ПП)\b', text, re.IGNORECASE)
            if alt:
                work_type = alt.group(1).upper()

    m5 = re.search(r'Номер:\s*(\d+)', text, re.IGNORECASE)
    if m5:
        num = m5.group(1)

    dept = ""
    m6 = re.search(r'Кафедра:\s*([^\n]+)', text, re.IGNORECASE)
    if m6:
        dept = m6.group(1).strip()

    if not work_type:
        work_type = get_work_type_from_text(text)
    if work_type == "ПР":
        work_type = "ПЗ"
    if work_type in ["КП", "КР"]:
        num = ""

    result = {
        "group": group, "author": author, "authors": authors,
        "subject": subject, "full_subject": subj_text,
        "work_type": work_type, "work_number": num,
        "department": dept, "raw_text": text,
    }
    result["filename"] = generate_safe_filename(result)
    return result

def extract_metadata(filename: str, text: str = "", is_presentation: bool = False) -> dict:
    base = os.path.basename(filename)
    work_type_from_name = None
    for wt in ALLOWED_WORK_TYPES:
        if f"_{wt}" in base or f"_{wt.lower()}" in base:
            work_type_from_name = wt
            break

    raw = normalize_text(clean_text_for_parsing(f"{filename} {text}"))
    raw_head = raw
    end_m = re.search(r'(?:содержание|введение|оглавление)', raw_head.lower())
    if end_m:
        raw_head = raw_head[:end_m.start()]
    if len(raw_head.strip()) < 50:
        raw_head = raw[:2000]
    raw_head_lower = raw_head.lower()

    work_type = None
    for key, val in {
        "курсовой проект": "КП", "курсовая проект": "КП", "курсовая работа": "КР",
        "курсовой работа": "КР", "курсов": "КР", "лабораторная работа": "ЛР",
        "лабораторная": "ЛР", "лаб. раб": "ЛР", "практическая работа": "ПЗ",
        "практическая": "ПЗ", "пр. раб": "ПЗ", "практ": "ПЗ",
        "индивидуальное задание": "ИЗ", "индивидуальная работа": "ИЗ",
        "самостоятельная работа": "СРС", "домашняя контрольная": "ДКР",
        "учебная практика": "УП", "производственная практика": "ПП",
    }.items():
        if key in raw_head_lower:
            work_type = val
            break
    if not work_type:
        alt = re.search(r'\b(ЛР|ПЗ|КР|КП|ИЗ|СРС|ДКР|СР|ПР|УП|ПП)\b', raw_head)
        if alt:
            work_type = alt.group(1).upper()
    if not work_type:
        work_type = get_work_type_from_text(raw_head)
    if work_type == "ПР":
        work_type = "ПЗ"

    if work_type_from_name and (not work_type or work_type == "РАБ"):
        work_type = work_type_from_name

    group = None
    grp = re.search(
        r'(?:гр\.?|группа|group)[\s:]*([А-ЯЁA-Zа-яёa-z]{2,4}[-\s]?\d{2,3}[А-ЯЁA-Zа-яёa-z0-9]?)',
        raw, re.IGNORECASE
    )
    if grp:
        group = grp.group(1).upper().replace(" ", "-")
        if len(group) >= 5 and group[-1] == '6' and group[-2].isdigit():
            group = group[:-1] + 'Б'

    authors = extract_authors(raw)
    author = authors[0] if authors else None
    if not author:
        am = re.search(
            r'(?:Выполнил|выполнил|Автор|автор|Студент|студент).{0,300}?'
            r'(?<![А-ЯЁа-яёa-zA-Z])([А-ЯЁ][а-яё]+)[\s,]+([А-ЯЁ])[\.\s,]*([А-ЯЁ])?',
            raw
        )
        if am:
            author = f"{am.group(1)} {am.group(2)}{am.group(3) or ''}".replace('.', '')

    subj_text = ""
    for marker in ["дисциплине", "предмет", "курсу", "на тему"]:
        idx = raw_head_lower.find(marker)
        if idx == -1:
            continue
        after = re.sub(r'^[\s:]+', '', raw_head[idx + len(marker):].strip())
        if after and after[0] in '«"':
            qm = re.search(r'^[«"]([^»"]+)[»"]', after)
            if qm:
                subj_text = qm.group(1).strip()
                break
        first_stop = len(after)
        for sw in ["лабораторная", "практическая", "курсовая", "курсовой",
                   "отчёт", "отчет", "работа", "проект", "вариант", "выполнил", "№", "студент"]:
            si = after.lower().find(sw)
            if si != -1 and si < first_stop:
                first_stop = si
        subj_text = after[:first_stop].strip()[:100]
        break

    subject = get_subject_abbreviation(subj_text) if subj_text else "Предмет"
    if not subj_text:
        subj_text = "Предмет"

    department = extract_department(raw_head)

    num = None
    sn = re.search(r'(?:работа|лр|пз|лабораторная|практическая)[\s:]*(?:№|номер)?[\s:]*(\d+)', raw_head, re.IGNORECASE)
    if sn:
        num = sn.group(1)
    else:
        nm = re.search(r'(?:№|номер)[\s:]*(\d+)', raw_head, re.IGNORECASE)
        if nm:
            num = nm.group(1)
        else:
            fn = re.search(r'(?:задание|вариант)[\s:]*(\d+)', raw_head, re.IGNORECASE)
            if fn:
                num = fn.group(1)

    if work_type in ["КП", "КР"]:
        num = ""

    result = {
        "group": group, "author": author, "authors": authors,
        "subject": subject, "full_subject": subj_text,
        "work_type": work_type, "work_number": num,
        "department": department, "raw_text": text,
    }

    # ---- БЛОК ДЛЯ ПРЕЗЕНТАЦИЙ ----
    if is_presentation:
        # Если авторы не были найдены – пробуем извлечь из списка участников
        if not result.get("authors"):
            raw_lower = raw.lower()
            # Ищем блок "Участники и роли" или "Участники:"
            parts = re.split(r'(?i)(?:участники\s+и\s+роли|участники:|команда:)', raw)
            if len(parts) > 1:
                block = parts[1]
                # Обрезаем до следующего заголовка (Введение, Цель и т.п.)
                block = re.split(r'(?i)(?:введение|цель|задачи|стек\s+технологий|макеты|диаграмма|тестирование)', block)[0]
                # Ищем строки вида "Фамилия Имя Отчество - роль" или "Фамилия И.О. - роль"
                lines = re.findall(r'([А-ЯЁ][а-яё]+\s+[А-ЯЁ][а-яё]+\s+[А-ЯЁ][а-яё]+)\s*[-—]\s*[^\n]+', block)
                if not lines:
                    # Пробуем вариант с инициалами: "Иванов А.О. - дизайнер"
                    lines = re.findall(r'([А-ЯЁ][а-яё]+\s+[А-ЯЁ]\.\s*[А-ЯЁ]\.)\s*[-—]\s*[^\n]+', block)
                if lines:
                    result["authors"] = [line.strip() for line in lines]
                    result["author"] = result["authors"][0]

        # Если предмет не найден или равен "Предмет" – пробуем взять тему из первого заголовка
        if not result.get("subject") or result.get("subject") == "Предмет":
            # Ищем первую строку, которая не является служебной
            lines = [line.strip() for line in raw.split('\n') if line.strip() and len(line.strip()) > 10]
            for line in lines:
                # Исключаем строки с "Участники", "Руководитель" и т.п.
                if re.search(r'(?i)(участники|руководитель|министерство|кафедра|утверждаю|отчёт|реферат|аннотация)', line):
                    continue
                # Если строка длинная и содержит ключевые слова (разработка, создание, платформа)
                if re.search(r'(?i)(разработк|создани|платформ|систем|сервис|приложени|веб|сайт)', line):
                    result["full_subject"] = line
                    # Сокращаем до 3-4 слов для аббревиатуры
                    words = line.split()
                    abbr = ''.join([w[0].upper() for w in words if len(w) > 2])[:10]
                    if abbr:
                        result["subject"] = abbr
                    else:
                        result["subject"] = "Проект"
                    break
            # Если не нашли – оставляем "Проектная деятельность"
            if not result.get("full_subject"):
                result["full_subject"] = "Проектная деятельность"
                result["subject"] = "Проект"

        # Если группа не найдена – можно извлечь из имени файла (если есть)
        if not result.get("group"):
            grp_from_name = re.search(r'([А-ЯЁA-Z]{2,4}[_\-\s]?\d{2,4}[А-ЯЁA-Z]?)', base, re.IGNORECASE)
            if grp_from_name:
                result["group"] = grp_from_name.group(1).upper().replace(" ", "-")

        # Если авторы всё ещё не найдены – пробуем извлечь из "Руководитель" (как fallback)
        if not result.get("authors"):
            leader_match = re.search(r'(?:руководитель|научный руководитель)\s*[:]?\s*([А-ЯЁ][а-яё]+\s+[А-ЯЁ]\.\s*[А-ЯЁ]\.)', raw, re.IGNORECASE)
            if leader_match:
                result["authors"] = [leader_match.group(1).strip()]
                result["author"] = leader_match.group(1).strip()

    result["filename"] = generate_safe_filename(result)
    return result



# ─── Презентации (задача 9) ──────────────────────────────

def extract_subject_from_pptx(pptx_path: str) -> str:
    try:
        prs = Presentation(pptx_path)
        if not prs.slides:
            return ""
        slide = prs.slides[0]
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
        full_text = "\n".join(lines)
        m = re.search(
            r'(?:по\s+дисциплине|по\s+предмету)[\s:«"]*([^\n»"\r]{3,120})',
            full_text, re.IGNORECASE
        )
        if m:
            subj = m.group(1).strip().strip('»"«"')
            subj = re.split(r'\n|Лектор|Преподаватель|Студент|Группа|Кафедра', subj, maxsplit=1)[0].strip()
            return subj[:150]
        return ""
    except Exception as e:
        logging.error(f"extract_subject_from_pptx error: {e}")
        return ""

# ─── Разбор титульного листа ──────────────────────────────

async def parse_title_page(file_path: str) -> dict:
    try:
        ext = os.path.splitext(file_path)[1].lower()
        text = ""
        logging.info(f"⚡ Разбор: {file_path}")

        if ext == '.docx':
            doc = docx.Document(file_path)
            lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_data = [c.text.strip() for c in row.cells if c.text.strip()]
                    if row_data:
                        lines.append(" ".join(row_data))
            text = re.sub(r'\s+', ' ', "\n".join(lines))
            logging.info(f"DOCX text[:400]: {text[:400]}")

            result = extract_metadata(file_path, text)
            need_gemini = (
                not result["group"] or result["group"] == "БЕЗ_ГРУППЫ" or
                not result["author"] or result["author"] == "Автор" or
                not result["subject"] or result["subject"] == "Предмет"
            )
            if need_gemini:
                result = await validate_metadata_with_gemini(result, text)
            result['is_ocr'] = False

        else:
            reader = PdfReader(file_path)
            text = reader.pages[0].extract_text() or ""
            cyrillic = len(re.findall(r'[а-яА-ЯёЁ]', text))

            result = extract_metadata(file_path, text)
            need_gemini = (
                not result["group"] or result["group"] == "БЕЗ_ГРУППЫ" or
                not result["author"] or result["author"] == "Автор" or
                not result["subject"] or result["subject"] == "Предмет"
            )
            if need_gemini or len(text.strip()) < 50 or cyrillic < 10:
                raw_ocr = await extract_raw_text_gemini(file_path)
                result = await validate_metadata_with_gemini(result, raw_ocr)
                result['is_ocr'] = True
            else:
                result['is_ocr'] = False

        if not result.get("filename") or result["filename"] == "error.pdf":
            result["filename"] = generate_safe_filename(result)
        return result

    except Exception as e:
        logging.error(f"parse_title_page error: {e}", exc_info=True)
        return {"filename": "error.pdf", "is_ocr": False}

# ─── Извлечение текста из PDF через fitz ─────────────────

async def validate_metadata_with_gemini(extracted_data: dict, raw_text: str) -> dict:
    logging.info("Валидация метаданных через Gemini...")
    api_key = os.environ.get("GOOGLE_API_KEY")
    if not api_key:
        return extracted_data
    ai_client = genai.Client(api_key=api_key)

    prompt = f"""
Ты сверяешь уже извлечённые поля документа с исходным OCR-текстом.
ДОПУСТИМЫЕ ТИПЫ РАБОТ: {", ".join(ALLOWED_WORK_TYPES)}

ЗАДАЧА:
1. Сверь текущие данные с текстом OCR. Пустые или явно неверные поля — исправь.
2. work_type: сопоставь с допустимым типом из списка.
3. group: перепиши шифр ровно как в тексте OCR, без подгонки к шаблону.
4. department: найди строку «Кафедра …» (КАФЕДРА, Кафедра, кафедра) и выпиши название полностью.
5. Нет данных — верни null. Не придумывай.

ТЕКУЩИЕ ДАННЫЕ:
{json.dumps(extracted_data, ensure_ascii=False)}

ТЕКСТ OCR:
{raw_text[:2000]}
"""
    schema = types.Schema(
        type=types.Type.OBJECT,
        properties={
            "group":       types.Schema(type=types.Type.STRING),
            "author":      types.Schema(type=types.Type.STRING),
            "subject":     types.Schema(type=types.Type.STRING),
            "work_type":   types.Schema(type=types.Type.STRING),
            "work_number": types.Schema(type=types.Type.STRING),
            "department":  types.Schema(type=types.Type.STRING),
        },
        required=["group", "author", "subject", "work_type", "work_number"]
    )
    try:
        response = await asyncio.wait_for(
            ai_client.aio.models.generate_content(
                model='gemini-3.1-flash-lite',
                contents=prompt,
                config=types.GenerateContentConfig(
                    temperature=0.0,
                    response_mime_type="application/json",
                    response_schema=schema,
                )
            ),
            timeout=30
        )
        validated = json.loads(response.text.strip())
        logging.info(f"Валидатор Gemini: {validated}")

        for key, value in validated.items():
            if value is None:
                continue
            if key == "author" and isinstance(value, (list, tuple)):
                value = next((v for v in value if v), None)
                if value is None:
                    continue
            # ВАЖНО: Gemini возвращает "сырое" значение из распознанного текста (особенно
            # при плохом OCR), без приведения к единому формату. Без нормализации здесь
            # group/author могут разойтись в написании с уже сохранёнными в базе версиями
            # того же документа ("ПЕ-41Б" vs "ПЕ41Б", "Сыропятов АВ" vs "Сыропятов А.В.") —
            # из-за этого find_exact_work перестаёт находить дубликат по точному совпадению
            # строк, и в базе накапливаются две записи одной и той же работы.
            if key == "group":
                value = clean_group_name(value)
            elif key == "author":
                value = normalize_author_name(value)
            extracted_data[key] = value

        if validated.get('subject'):
            extracted_data['subject'] = get_subject_abbreviation(validated['subject'])
            extracted_data['full_subject'] = validated['subject']
        elif extracted_data.get('full_subject'):
            extracted_data['subject'] = get_subject_abbreviation(extracted_data['full_subject'])

        va = validated.get('author')
        if isinstance(va, (list, tuple)):
            va = next((v for v in va if v), None)
        if va and not extracted_data.get('authors'):
            extracted_data['authors'] = [va]

        # !!! ВАЖНО: приводим кафедру к короткому коду
        if validated.get('department'):
            extracted_data['department'] = extract_department(validated['department'])
        elif extracted_data.get('department'):
            extracted_data['department'] = extract_department(extracted_data['department'])

        extracted_data["filename"] = generate_safe_filename(extracted_data)
        return extracted_data
    except Exception as e:
        logging.error(f"validate_metadata_with_gemini error: {e}")
        return extracted_data
# ----- Функции для определения подписи и сканированных страниц -----
